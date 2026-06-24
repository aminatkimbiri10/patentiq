#!/usr/bin/env python3
"""Génère le PowerPoint de soutenance PatentIQ — I2PA (scripts oraux inclus)."""

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUTPUT = ROOT / "docs" / "SOUTENANCE_PATENTIQ.pptx"
SCRIPTS_MD = ROOT / "docs" / "SOUTENANCE_SCRIPTS_ORAUX.md"

BRAND = RGBColor(10, 102, 194)  # #0A66C2
DARK = RGBColor(30, 41, 59)
MUTED = RGBColor(100, 116, 139)
WHITE = RGBColor(255, 255, 255)

SLIDES = [
    {
        "layout": "title",
        "title": "PatentIQ",
        "subtitle": "Plateforme web d'assistance à la propriété intellectuelle",
        "footer": "[Prénom NOM] · Stage I2PA · [JJ/MM/AAAA – JJ/MM/AAAA]\n[École / Filière] · Encadrante : [Nom] · Tuteur : [Nom]",
    },
    {
        "title": "Plan de la présentation",
        "bullets": [
            "Contexte I2PA et problématique",
            "Objectifs et méthode de travail",
            "Architecture et acteurs",
            "Les 4 piliers fonctionnels de PatentIQ",
            "Intelligence artificielle (encadrée et honnête)",
            "Résultats, difficultés et perspectives",
        ],
    },
    {
        "title": "I2PA — Contexte d'accueil",
        "bullets": [
            "Cabinet marocain de conseil en propriété intellectuelle (Mohammedia)",
            "Mission : accompagner porteurs et entreprises — brevets, marques, dessins",
            "I2PA conseille et structure ; le dépôt officiel reste sur l'OMPIC (directompic.ma)",
            "Slogan : « Protéger, valoriser, innover »",
        ],
    },
    {
        "title": "Problématique",
        "bullets": [
            "Dossiers PI dispersés : emails, PDF mal rangés, antériorité non reliée au dossier",
            "Délais critiques : opposition marque (~2 mois) vs publication brevet (~18 mois)",
            "Pas d'outil léger adapté au contexte marocain porteur ↔ conseiller CPI",
        ],
        "quote": "Comment centraliser et fiabiliser le suivi d'un dossier PI, sans se substituer à l'OMPIC, avec une IA transparente ?",
    },
    {
        "title": "Objectifs du stage",
        "bullets": [
            "Comprendre le parcours réel d'un dossier PI au Maroc",
            "Concevoir et développer PatentIQ par itérations courtes",
            "Aligner les fonctionnalités sur les retours de l'encadrante",
            "Tester, documenter et préparer une démonstration crédible",
        ],
    },
    {
        "title": "Méthode de travail",
        "bullets": [
            "Développement incrémental : migration SQL → action serveur → interface",
            "Validation « double session » : porteur (navigateur 1) + CPI (navigateur 2)",
            "3 phases : socle collaboratif → attentes métier → consolidation",
            "32 migrations · 104+ tests unitaires · documentation dans docs/",
        ],
    },
    {
        "title": "Architecture technique",
        "bullets": [
            "Frontend : Next.js 14 (App Router), React, TypeScript, Tailwind, shadcn/ui",
            "Backend : Server Actions + Supabase (PostgreSQL, Auth, Storage)",
            "Sécurité : Row Level Security — permissions en base, pas seulement dans l'UI",
            "Externe : EPO OPS (brevets), Hugging Face (synthèse), search.ompic.ma (marques)",
            "Déploiement : Vercel + Supabase (tier gratuit)",
        ],
    },
    {
        "title": "Acteurs et espaces",
        "bullets": [
            "Porteur de projet → /dashboard : crée le dossier, dépose, lance les analyses",
            "Conseiller CPI → /cpi : revue, statuts, Kanban, surveillance portefeuille",
            "Expert métier → /expert : avis technique structuré",
            "Administrateur → /admin : utilisateurs, paramètres, audit",
            "7 catégories de dossiers · checklists adaptées à chaque type de titre",
        ],
    },
    {
        "title": "Pilier 1 — Dossier collaboratif",
        "bullets": [
            "Projet structuré : Dossier · Analyses IA · Échanges",
            "Documents versionnés (PDF, Word, images) — stockage sécurisé",
            "Checklist métier par catégorie (brevet ≠ marque ≠ dessin)",
            "Score de complétude : résumé, documents, checklist → % + prochaine action",
            "Workflow 9 statuts : brouillon → soumis → revue CPI → approuvé / clôturé",
        ],
    },
    {
        "title": "Pilier 2 — Rédaction & revendications",
        "bullets": [
            "Sections brevet structurées (titre, abrégé, description, revendications)",
            "Revendications dans un espace dédié — séparé des messages ouverts",
            "Pré-examen automatique : anomalies fréquentes avant revue CPI",
            "  · Revendication indépendante, abrégé ≤ 150 mots, termes vagues, base d'antériorité",
            "Export dossier brevet (PDF) pour le conseiller",
        ],
    },
    {
        "title": "Pilier 3 — Analyses IA assistées",
        "bullets": [
            "7 types : nouveauté, FTO, sémantique, similarité, résumé, classification IPC, tags",
            "Sources réelles : EPO OPS + filtre pn=MA (brevets Maroc) + OMPIC marques",
            "Pipeline asynchrone : pending → worker → completed (documenté)",
            "Synthèses indicatives — rappel explicite : ne remplace pas un avis juridique",
            "Incident corrigé : pa=MA (déposant) → pn=MA (publication marocaine)",
        ],
    },
    {
        "title": "Pilier 4 — Surveillance & veille",
        "bullets": [
            "Watchlist OMPIC : scan similarité marques sur search.ompic.ma",
            "Modes honnêtes : live · stub (démo) · proxy · hybrid",
            "Veille technologique : mots-clés + classes IPC sur brevets EPO",
            "Échéances PI : opposition marque, publication brevet — rappels tableau de bord",
            "Dessins : cycle de gestion oui ; recherche auto non (pas d'API OMPIC)",
        ],
    },
    {
        "title": "Cycles de vie distincts",
        "bullets": [
            "Marque : dépôt → publication → opposition (~2 mois) → enregistrement → surveillance",
            "Brevet : dépôt → examen → publication (~18 mois) → délivrance → veille",
            "Dessin & modèle : dépôt → examen → publication → surveillance manuelle",
            "Libellés UI adaptés : « Description de la marque » ≠ « Résumé de l'invention »",
        ],
    },
    {
        "title": "Fonctionnalités avancées (fin de stage)",
        "bullets": [
            "Tableau de bord portefeuille CPI (répartition statuts / types PI)",
            "Cartographie des revendications vs antériorité (claim chart)",
            "Générateur de dénominations de marque + vérification OMPIC",
            "Assistance réponse aux irrégularités OMPIC (brouillon structuré)",
            "Pré-examen brevet + score de complétude (heuristiques testables)",
        ],
    },
    {
        "title": "Sécurité & confidentialité",
        "bullets": [
            "PostgreSQL RLS : can_view_project(), can_edit_project()…",
            "Un porteur ne voit jamais le dossier d'un autre",
            "Email confirmé à l'inscription · 2FA TOTP disponible",
            "Journal d'audit des actions sensibles (admin)",
        ],
    },
    {
        "title": "Tests & validation",
        "bullets": [
            "104+ tests unitaires (Vitest) : EPO, cycles PI, complétude, pré-examen…",
            "Tests E2E Playwright : parcours critiques",
            "Validation métier : démos régulières devant l'encadrante",
            "Scénarios documentés : parcours marque + parcours brevet",
        ],
    },
    {
        "title": "Résultats à la clôture",
        "bullets": [
            "✓ Parcours porteur → conseiller opérationnel",
            "✓ Analyses IA + surveillance OMPIC (mode hybrid)",
            "✓ Rédaction brevet, revendications, export",
            "✓ Cycles marque / brevet / dessin différenciés",
            "✓ Déploiement documenté (Vercel + Supabase)",
        ],
    },
    {
        "title": "Limites assumées",
        "bullets": [
            "Pas d'API OMPIC officielle → recherche marques via portail public",
            "Analyses IA : worker requis (npm run ai:worker:loop)",
            "Synthèses indicatives — décision finale = CPI",
            "Dessins : pas de moteur de recherche automatique (choix de transparence)",
            "MVP crédible, pas produit commercial complet",
        ],
    },
    {
        "title": "Perspectives",
        "bullets": [
            "Industrialiser la passerelle OMPIC (partenariat / proxy)",
            "Notifications email · embeddings vectoriels pour le RAG",
            "Renforcement tests E2E en CI",
            "Mise en production encadrée chez I2PA",
        ],
    },
    {
        "title": "Conclusion",
        "bullets": [
            "PatentIQ structure la préparation et le suivi d'un dossier PI pour I2PA",
            "Collaboration multi-acteurs · antériorité assistée · surveillance OMPIC",
            "Le métier impose l'architecture — écouter avant de coder",
            "Le dépôt officiel reste sur directompic.ma",
        ],
        "quote": "« Une limite documentée vaut mieux qu'une illusion confortable. »",
    },
    {
        "layout": "title",
        "title": "Merci pour votre attention",
        "subtitle": "Questions ?",
        "footer": "PatentIQ · I2PA · [Prénom NOM]\nDémo live disponible · docs/DEMO_ENCADRANTE.md",
    },
]


def load_oral_scripts() -> list[str]:
    """Lit les scripts oraux depuis docs/SOUTENANCE_SCRIPTS_ORAUX.md (21 slides)."""
    if not SCRIPTS_MD.exists():
        return []
    text = SCRIPTS_MD.read_text(encoding="utf-8")
    parts = re.findall(
        r"## Slide \d+[^\n]*\n(.*?)(?=\n---\n|\n## Banque|\Z)",
        text,
        flags=re.DOTALL,
    )
    return [p.strip() for p in parts if p.strip()]


def set_paragraph_style(paragraph, size=18, bold=False, color=DARK, align=None):
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    if align:
        paragraph.alignment = align


def add_slide_notes(slide, oral_script: str | None, slide_num: int):
    if not oral_script:
        return
    header = f"SCRIPT ORAL — Slide {slide_num}\n{'─' * 40}\n\n"
    slide.notes_slide.notes_text_frame.text = header + oral_script


def add_title_slide(prs, data, oral_script: str | None, slide_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # bandeau haut
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(2.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRAND
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = data["title"]
    set_paragraph_style(p, size=40, bold=True, color=WHITE)

    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11), Inches(1.5))
    tf2 = sub_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = data.get("subtitle", "")
    set_paragraph_style(p2, size=22, color=DARK)

    if data.get("footer"):
        foot = slide.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(11), Inches(1.2))
        tf3 = foot.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = data["footer"]
        set_paragraph_style(p3, size=14, color=MUTED)

    add_slide_notes(slide, oral_script, slide_num)


def add_content_slide(prs, data, oral_script: str | None, slide_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # barre titre
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.12), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = BRAND
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.8))
    p = title_box.text_frame.paragraphs[0]
    p.text = data["title"]
    set_paragraph_style(p, size=28, bold=True, color=BRAND)

    body = slide.shapes.add_textbox(Inches(0.55), Inches(1.25), Inches(12), Inches(4.8))
    tf = body.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(data.get("bullets", [])):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        text = bullet.strip()
        if text.startswith("  ·"):
            para.text = text.replace("  ·", "•", 1)
            para.level = 1
            set_paragraph_style(para, size=16, color=MUTED)
        else:
            para.text = text
            para.level = 0
            set_paragraph_style(para, size=18, color=DARK)
        para.space_after = Pt(8)

    if data.get("quote"):
        qbox = slide.shapes.add_textbox(Inches(0.55), Inches(5.9), Inches(12), Inches(1))
        qp = qbox.text_frame.paragraphs[0]
        qp.text = data["quote"]
        set_paragraph_style(qp, size=15, color=BRAND, align=PP_ALIGN.CENTER)
        qp.font.italic = True

    # numéro slide footer
    footer = slide.shapes.add_textbox(Inches(11), Inches(7.1), Inches(2), Inches(0.3))
    fp = footer.text_frame.paragraphs[0]
    fp.text = "PatentIQ · I2PA"
    set_paragraph_style(fp, size=10, color=MUTED, align=PP_ALIGN.RIGHT)

    add_slide_notes(slide, oral_script, slide_num)


def main():
    oral_scripts = load_oral_scripts()
    if len(oral_scripts) != len(SLIDES):
        print(
            f"Attention : {len(oral_scripts)} scripts lus, {len(SLIDES)} slides — "
            "vérifiez docs/SOUTENANCE_SCRIPTS_ORAUX.md"
        )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for i, data in enumerate(SLIDES):
        script = oral_scripts[i] if i < len(oral_scripts) else None
        slide_num = i + 1
        if data.get("layout") == "title":
            add_title_slide(prs, data, script, slide_num)
        else:
            add_content_slide(prs, data, script, slide_num)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Présentation générée : {OUTPUT}")
    print(f"Slides : {len(SLIDES)} · Scripts oraux : {len(oral_scripts)}")


if __name__ == "__main__":
    main()
